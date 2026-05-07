import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "TenderMind_AI_Extended_Code_Walkthrough.pptx"
SOURCE_PDF = Path(r"C:\Users\Shivansh Karanwal\Downloads\Team Takshak - AI for Bharat.pdf")
ASSET_DIR = ROOT / "data" / "ppt_assets"

NAVY = RGBColor(18, 32, 51)
BLUE = RGBColor(39, 96, 172)
TEAL = RGBColor(27, 132, 121)
GREEN = RGBColor(29, 116, 72)
RED = RGBColor(170, 52, 45)
AMBER = RGBColor(184, 123, 22)
PURPLE = RGBColor(102, 68, 167)
SLATE = RGBColor(74, 85, 104)
LIGHT_BG = RGBColor(247, 249, 252)
LIGHT_BLUE = RGBColor(228, 238, 252)
LIGHT_GREEN = RGBColor(221, 247, 231)
LIGHT_RED = RGBColor(255, 225, 223)
LIGHT_AMBER = RGBColor(255, 242, 202)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(20, 24, 34)
BORDER = RGBColor(216, 222, 232)


def main() -> None:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    render_source_pdf_cover()
    result = load_latest_result()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)
    add_problem_slide(prs)
    add_deck_to_code_slide(prs)
    add_demo_story_slide(prs)
    add_architecture_slide(prs)
    add_repo_map_slide(prs)
    add_api_slide(prs)
    add_models_slide(prs)
    add_stage1_slide(prs)
    add_stage2_slide(prs)
    add_field_extraction_slide(prs)
    add_stage3_slide(prs)
    add_rule_engine_slide(prs)
    add_semantic_slide(prs)
    add_confidence_slide(prs)
    add_hitl_slide(prs)
    add_reporting_audit_slide(prs)
    add_sample_result_slide(prs, result)
    add_explainability_slide(prs, result)
    add_demo_run_slide(prs)
    add_roadmap_slide(prs)
    add_closing_slide(prs)

    prs.save(OUTPUT)
    print(f"Created {OUTPUT}")


def render_source_pdf_cover() -> None:
    if not SOURCE_PDF.exists():
        return
    cover = ASSET_DIR / "source_deck_cover.png"
    if cover.exists():
        return
    try:
        import fitz

        doc = fitz.open(str(SOURCE_PDF))
        page = doc[0]
        pix = page.get_pixmap(dpi=120)
        pix.save(str(cover))
    except Exception:
        pass


def load_latest_result():
    report_dir = ROOT / "data" / "reports"
    reports = sorted(report_dir.glob("eval_*.json"), key=lambda path: path.stat().st_mtime)
    if not reports:
        return None
    return json.loads(reports[-1].read_text(encoding="utf-8"))


def blank_slide(prs, title, kicker=None, dark=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, dark=dark)
    if kicker:
        add_text(slide, kicker.upper(), 0.6, 0.35, 6.0, 0.24, Pt(9), TEAL if not dark else RGBColor(130, 220, 210), bold=True)
    add_text(slide, title, 0.58, 0.62, 10.6, 0.55, Pt(27), WHITE if dark else NAVY, bold=True)
    add_footer(slide, dark=dark)
    return slide


def add_background(slide, dark=False):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY if dark else LIGHT_BG
    if not dark:
        top = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.18))
        top.fill.solid()
        top.fill.fore_color.rgb = TEAL
        top.line.fill.background()


def add_footer(slide, dark=False):
    color = RGBColor(184, 197, 214) if dark else RGBColor(116, 126, 145)
    add_text(slide, "TenderMind AI | AI-based tender evaluation system", 0.6, 7.12, 6.8, 0.22, Pt(8), color)


def add_text(slide, text, x, y, w, h, size=Pt(14), color=BLACK, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Aptos"
    if align:
        p.alignment = align
    return box


def add_bullets(slide, items, x, y, w, h, size=Pt(14), color=BLACK, gap=0.05):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"- {item}"
        p.level = 0
        p.font.size = size
        p.font.color.rgb = color
        p.font.name = "Aptos"
        p.space_after = Pt(6 + gap)
    return box


def card(slide, title, body, x, y, w, h, fill=WHITE, accent=TEAL, title_size=Pt(14), body_size=Pt(11)):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(1)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    add_text(slide, title, x + 0.18, y + 0.15, w - 0.35, 0.28, title_size, NAVY, bold=True)
    if isinstance(body, list):
        add_bullets(slide, body, x + 0.22, y + 0.55, w - 0.36, h - 0.65, body_size, BLACK)
    else:
        add_text(slide, body, x + 0.18, y + 0.55, w - 0.34, h - 0.65, body_size, BLACK)
    return shape


def pill(slide, text, x, y, w, fill, color=WHITE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.34))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    add_text(slide, text, x, y + 0.075, w, 0.16, Pt(9), color, bold=True, align=PP_ALIGN.CENTER)
    return shape


def connector(slide, x1, y1, x2, y2, color=SLATE):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(1.5)
    return line


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, dark=True)
    add_text(slide, "TenderMind AI", 0.65, 0.65, 6.5, 0.55, Pt(34), WHITE, bold=True)
    add_text(slide, "Extended technical walkthrough of the working prototype", 0.68, 1.25, 6.7, 0.35, Pt(15), RGBColor(203, 214, 229))
    add_text(slide, "From tender PDF upload to explainable PASS / FAIL / NEEDS_REVIEW reports.", 0.68, 1.75, 6.4, 0.48, Pt(18), WHITE)

    card(slide, "Built for demo and review", [
        "FastAPI backend and Streamlit frontend",
        "Four-stage AI pipeline with audit logging",
        "No API key required for local hackathon demo",
        "Human review for ambiguity instead of silent rejection",
    ], 0.72, 2.55, 5.7, 2.45, fill=RGBColor(238, 243, 250), accent=TEAL)

    cover = ASSET_DIR / "source_deck_cover.png"
    if cover.exists():
        slide.shapes.add_picture(str(cover), Inches(7.35), Inches(0.75), width=Inches(4.8))
        add_text(slide, "Source deck context used: Team Takshak - AI for Bharat", 7.35, 6.04, 4.9, 0.3, Pt(9), RGBColor(203, 214, 229))
    else:
        card(slide, "Source context", "Team Takshak pitch deck plus the implemented codebase.", 7.2, 1.25, 4.6, 2.4, fill=RGBColor(238, 243, 250), accent=BLUE)
    add_footer(slide, dark=True)


def add_problem_slide(prs):
    slide = blank_slide(prs, "Why tender evaluation needs structure", "Problem")
    card(slide, "Manual process today", [
        "Tender criteria are scattered across long legal documents.",
        "Bidder submissions mix typed PDFs, scans, images, and inconsistent formats.",
        "Officers build informal checklists, so consistency depends on memory and time.",
    ], 0.7, 1.45, 3.8, 2.35, fill=WHITE, accent=RED)
    card(slide, "Operational failures", [
        "Eligible bidder can be rejected because a scan is misread.",
        "Ineligible bidder can pass when a certificate or declaration is missed.",
        "Committee cannot always defend the reasoning later.",
    ], 4.75, 1.45, 3.8, 2.35, fill=WHITE, accent=AMBER)
    card(slide, "TenderMind fix", [
        "Extract criteria into a structured schema.",
        "Map bidder evidence to fields with confidence scores.",
        "Evaluate criterion by criterion with explainable rules.",
        "Flag uncertainty for human review.",
    ], 8.8, 1.45, 3.8, 2.35, fill=WHITE, accent=TEAL)
    add_text(slide, "Design promise", 0.75, 4.35, 2.3, 0.3, Pt(14), NAVY, bold=True)
    for idx, text in enumerate(["Zero silent disqualification", "Criterion-level reasoning", "Auditable event trail", "Officer remains final authority"]):
        pill(slide, text, 0.75 + idx * 3.0, 4.85, 2.55, [TEAL, BLUE, PURPLE, GREEN][idx])


def add_deck_to_code_slide(prs):
    slide = blank_slide(prs, "How the pitch became a runnable prototype", "Bridge")
    card(slide, "Deck concept", [
        "4-stage AI pipeline",
        "Explainable verdicts",
        "Human-in-the-loop ambiguity handling",
        "Audit trail and compliance positioning",
    ], 0.85, 1.55, 3.2, 3.45, fill=WHITE, accent=BLUE)
    card(slide, "Implemented codebase", [
        "FastAPI endpoints for upload, extraction, evaluation, report",
        "Pydantic models for criteria, bidders, fields, verdicts",
        "PDF/OCR document parser and ZIP bidder handling",
        "Rule engine, semantic matcher, confidence logic",
        "SQLite audit logs plus JSON and HTML reports",
    ], 5.0, 1.2, 3.45, 4.15, fill=WHITE, accent=TEAL)
    card(slide, "Demo-ready result", [
        "Runs locally in minutes",
        "Uses mock LLM fallback by default",
        "Sample tender and 3 bidders included",
        "Color-coded Streamlit interface",
    ], 9.4, 1.55, 3.0, 3.45, fill=WHITE, accent=GREEN)
    connector(slide, 4.1, 3.1, 4.9, 3.1, BLUE)
    connector(slide, 8.5, 3.1, 9.3, 3.1, BLUE)


def add_demo_story_slide(prs):
    slide = blank_slide(prs, "End-to-end working flow", "Demo story")
    stages = [
        ("1", "Upload tender PDF", "Tender saved under data/uploads/tenders/<id>"),
        ("2", "Extract criteria", "LLM abstraction returns validated Criterion objects"),
        ("3", "Upload bidder ZIP", "ZIP safely extracted and grouped by bidder"),
        ("4", "Parse documents", "PDF text first; OCR fallback for scans/images"),
        ("5", "Evaluate", "Rules + semantic similarity + confidence thresholds"),
        ("6", "Report", "Matrix, explanations, JSON/HTML, audit log"),
    ]
    x = 0.65
    for idx, (num, title, body) in enumerate(stages):
        y = 1.35 + (idx % 3) * 1.65
        col_x = x + (idx // 3) * 6.15
        pill(slide, f"Stage {num}", col_x, y, 1.0, TEAL if idx < 3 else BLUE)
        add_text(slide, title, col_x + 1.15, y - 0.01, 4.4, 0.28, Pt(14), NAVY, bold=True)
        add_text(slide, body, col_x + 1.15, y + 0.37, 4.8, 0.45, Pt(10.5), SLATE)
        if idx % 3 != 2:
            connector(slide, col_x + 0.5, y + 0.45, col_x + 0.5, y + 1.35, BORDER)
    connector(slide, 6.0, 3.45, 6.65, 3.45, TEAL)


def add_architecture_slide(prs):
    slide = blank_slide(prs, "System architecture overview", "Architecture")
    card(slide, "Frontend", ["Streamlit app", "Tender PDF upload", "Bidder ZIP upload", "Matrix and flagged cases"], 0.65, 1.25, 2.55, 2.35, fill=WHITE, accent=BLUE)
    card(slide, "FastAPI", ["/upload-tender", "/upload-bidders", "/extract-criteria", "/evaluate", "/report"], 3.55, 1.25, 2.55, 2.35, fill=WHITE, accent=TEAL)
    card(slide, "AI services", ["CriteriaExtractor", "DocumentParser", "MatchingEngine", "ReportService"], 6.45, 1.25, 2.55, 2.35, fill=WHITE, accent=PURPLE)
    card(slide, "Storage", ["data/uploads", "data/criteria", "data/reports", "SQLite audit DB"], 9.35, 1.25, 2.55, 2.35, fill=WHITE, accent=GREEN)
    for x in [3.22, 6.12, 9.02]:
        connector(slide, x, 2.42, x + 0.28, 2.42, SLATE)

    add_text(slide, "Key runtime separation", 0.75, 4.35, 3.8, 0.28, Pt(14), NAVY, bold=True)
    add_bullets(slide, [
        "Frontend never evaluates directly; it calls API endpoints.",
        "Backend services isolate extraction, parsing, matching, and reporting.",
        "Pydantic schemas keep data contracts explicit between stages.",
        "Audit logging records important events with timestamps and metadata.",
    ], 0.85, 4.75, 8.4, 1.45, Pt(12))
    card(slide, "Deployment posture", "Local-first prototype today; same service boundaries can move to cloud storage, managed DB, queue workers, and authenticated officer workflows later.", 9.55, 4.55, 2.65, 1.55, fill=LIGHT_BLUE, accent=BLUE)


def add_repo_map_slide(prs):
    slide = blank_slide(prs, "Repository map: where each responsibility lives", "Codebase")
    columns = [
        ("backend/main.py", "FastAPI app and endpoints"),
        ("backend/models/schemas.py", "Pydantic contracts"),
        ("backend/services/", "Pipeline implementation"),
        ("backend/utils/", "Audit, file, LLM helpers"),
        ("frontend/streamlit_app.py", "Demo UI"),
        ("scripts/", "Sample data and smoke test"),
        ("data/sample/", "Mock tender and bidders"),
        ("data/reports/", "Generated JSON/HTML reports"),
    ]
    for idx, (path, role) in enumerate(columns):
        col = idx % 2
        row = idx // 2
        x = 0.75 + col * 6.1
        y = 1.2 + row * 1.25
        card(slide, path, role, x, y, 5.55, 0.85, fill=WHITE, accent=TEAL if col == 0 else BLUE, title_size=Pt(12), body_size=Pt(10.5))


def add_api_slide(prs):
    slide = blank_slide(prs, "API surface and request sequence", "Backend")
    endpoints = [
        ("POST /upload-tender", "Stores PDF and logs upload"),
        ("POST /upload-bidders", "Stores ZIP and logs upload"),
        ("POST /extract-criteria", "Creates criteria JSON for tender"),
        ("POST /evaluate", "Parses bidders and runs evaluation"),
        ("GET /report", "Returns latest or selected JSON/HTML report"),
        ("GET /audit-logs", "Returns timestamped audit events"),
    ]
    for idx, (endpoint, desc) in enumerate(endpoints):
        y = 1.25 + idx * 0.76
        pill(slide, endpoint, 0.75, y, 2.5, BLUE if idx < 4 else GREEN)
        add_text(slide, desc, 3.55, y + 0.05, 5.2, 0.2, Pt(11), BLACK)
    card(slide, "Endpoint design idea", [
        "Uploads create stable IDs.",
        "Extraction can be run separately so officers can inspect criteria first.",
        "Evaluation reuses saved criteria if available.",
        "Reports are deterministic artifacts for review and sharing.",
    ], 9.0, 1.25, 3.35, 3.55, fill=WHITE, accent=TEAL)


def add_models_slide(prs):
    slide = blank_slide(prs, "Data contracts: the backbone of explainability", "Pydantic schemas")
    card(slide, "Criterion", [
        "name, description, field_name",
        "data_type, operator, threshold",
        "mandatory, weight, evidence_keywords",
        "source_excerpt from tender",
    ], 0.75, 1.2, 3.45, 2.1, fill=WHITE, accent=BLUE)
    card(slide, "Bidder + ExtractedField", [
        "bidder id and name",
        "documents parsed from ZIP",
        "field values with confidence",
        "source document and excerpt",
    ], 4.75, 1.2, 3.45, 2.1, fill=WHITE, accent=TEAL)
    card(slide, "CriterionEvaluation", [
        "decision: PASS / FAIL / NEEDS_REVIEW",
        "confidence score",
        "extracted value",
        "rule applied",
        "reasoning and flagged marker",
    ], 8.75, 1.2, 3.45, 2.1, fill=WHITE, accent=PURPLE)
    add_text(slide, "Why this matters", 0.85, 4.05, 2.4, 0.3, Pt(15), NAVY, bold=True)
    add_bullets(slide, [
        "Every stage receives and returns typed objects instead of loose dictionaries.",
        "The report can show exactly which criterion, field, rule, confidence, and source produced a verdict.",
        "Future LLM/API integrations can change internally while keeping the schema stable.",
    ], 0.95, 4.45, 10.8, 1.2, Pt(12))


def add_stage1_slide(prs):
    slide = blank_slide(prs, "Stage 1: Criteria extraction from tender PDF", "Pipeline")
    card(slide, "Input", ["Tender PDF", "May contain clauses across sections", "Typed PDF supported now"], 0.7, 1.3, 2.7, 2.15, fill=WHITE, accent=BLUE)
    card(slide, "Text extraction", ["DocumentParser.extract_text_from_path", "pdfplumber first", "PyMuPDF fallback", "OCR fallback if needed"], 3.9, 1.3, 2.7, 2.15, fill=WHITE, accent=TEAL)
    card(slide, "LLM abstraction", ["LLMClient.extract_criteria", "Mock extractor by default", "OpenAI-ready wrapper point", "Returns Pydantic Criterion list"], 7.1, 1.3, 2.7, 2.15, fill=WHITE, accent=PURPLE)
    card(slide, "Output", ["data/criteria/<tender_id>.json", "Criteria visible before evaluation", "Officer can validate next"], 10.3, 1.3, 2.25, 2.15, fill=WHITE, accent=GREEN)
    for x in [3.45, 6.65, 9.85]:
        connector(slide, x, 2.38, x + 0.35, 2.38, SLATE)
    card(slide, "Prototype extraction examples", [
        "Minimum years of experience -> years_experience >= 3",
        "Average annual turnover -> annual_turnover_lakhs >= 50",
        "GST registration -> gst_registration is_true",
        "No blacklisting -> blacklisted_status is_false",
    ], 0.95, 4.25, 11.0, 1.35, fill=LIGHT_BLUE, accent=BLUE)


def add_stage2_slide(prs):
    slide = blank_slide(prs, "Stage 2: Document processing for bidder ZIPs", "Pipeline")
    card(slide, "ZIP handling", [
        "safe_extract_zip blocks unsafe paths",
        "Files grouped by bidder folder",
        "Supports PDF, image, and text extensions",
    ], 0.75, 1.25, 3.3, 2.1, fill=WHITE, accent=BLUE)
    card(slide, "PDF path", [
        "Try pdfplumber for typed PDFs",
        "Fallback to PyMuPDF text extraction",
        "Confidence around 0.90-0.93 for typed text",
    ], 4.45, 1.25, 3.3, 2.1, fill=WHITE, accent=TEAL)
    card(slide, "OCR path", [
        "Render scanned PDFs with PyMuPDF",
        "OpenCV grayscale + median blur + Otsu threshold",
        "pytesseract extracts image text",
    ], 8.15, 1.25, 3.3, 2.1, fill=WHITE, accent=PURPLE)
    card(slide, "Output object", [
        "DocumentExtraction(file_name, parser, text, confidence)",
        "Bidder(name, documents, fields)",
        "Each field stores source document and excerpt",
    ], 2.0, 4.15, 9.1, 1.35, fill=LIGHT_GREEN, accent=GREEN)


def add_field_extraction_slide(prs):
    slide = blank_slide(prs, "Field extraction: turning evidence into structured values", "Document parser")
    fields = [
        ("years_experience", "Regex: years / experience -> float"),
        ("annual_turnover_lakhs", "Crore/lakh normalization -> INR lakh"),
        ("similar_projects_count", "Project count -> integer"),
        ("gst_registration", "GSTIN or registered text -> boolean"),
        ("iso_9001_certificate", "valid/available/not available -> boolean"),
        ("blacklisted_status", "blacklisted/debarred declaration -> boolean"),
        ("net_worth_positive", "positive/negative net worth -> boolean"),
    ]
    for idx, (field, desc) in enumerate(fields):
        y = 1.15 + idx * 0.68
        pill(slide, field, 0.75, y, 3.1, TEAL)
        add_text(slide, desc, 4.15, y + 0.05, 5.8, 0.22, Pt(10.7), BLACK)
    card(slide, "Confidence metadata", [
        "Parser-level confidence starts at extraction.",
        "Exact label matches get a small boost.",
        "Every field carries raw value, normalized value, source file, and excerpt.",
        "Missing field is not a hard fail by itself; it becomes review.",
    ], 10.05, 1.2, 2.6, 4.55, fill=WHITE, accent=AMBER)


def add_stage3_slide(prs):
    slide = blank_slide(prs, "Stage 3: Matching engine", "Evaluation")
    card(slide, "Loop", [
        "For each bidder",
        "For each criterion",
        "Find matching extracted field",
        "Choose rule or semantic path",
    ], 0.75, 1.35, 3.1, 2.5, fill=WHITE, accent=BLUE)
    card(slide, "Rule engine", [
        "Numeric comparison",
        "Boolean truth checks",
        "Date comparison support",
        "Contains / between support",
    ], 4.1, 1.35, 3.1, 2.5, fill=WHITE, accent=TEAL)
    card(slide, "Semantic engine", [
        "MiniLM sentence-transformers if installed",
        "Cosine similarity with normalized embeddings",
        "Token/sequence fallback for local demo",
    ], 7.45, 1.35, 3.1, 2.5, fill=WHITE, accent=PURPLE)
    card(slide, "Result", [
        "CriterionEvaluation object",
        "decision + confidence",
        "rule_applied + reasoning",
        "flagged if uncertain",
    ], 10.8, 1.35, 1.75, 2.5, fill=WHITE, accent=GREEN)
    add_text(slide, "Important behavior: the engine evaluates every bidder against every criterion, creating a complete matrix instead of only reporting failures.", 0.95, 4.65, 11.3, 0.6, Pt(16), NAVY, bold=True)


def add_rule_engine_slide(prs):
    slide = blank_slide(prs, "Rule engine: deterministic, auditable comparisons", "Matching logic")
    card(slide, "Numeric", [
        ">=, <=, >, <, ==, between",
        "Examples: experience >= 3, turnover >= 50 lakh",
        "Values normalized before comparison",
    ], 0.75, 1.2, 3.45, 2.25, fill=WHITE, accent=BLUE)
    card(slide, "Boolean", [
        "is_true / is_false",
        "GST registration, ISO certificate, net worth",
        "No blacklisting uses is_false",
    ], 4.75, 1.2, 3.45, 2.25, fill=WHITE, accent=TEAL)
    card(slide, "Dates and text", [
        "date_before / date_after variants",
        "contains for phrase checks",
        "Unsupported or unsafe parsing -> review",
    ], 8.75, 1.2, 3.45, 2.25, fill=WHITE, accent=PURPLE)
    card(slide, "Audit-friendly output", [
        "matched: true/false",
        "confidence",
        "rule_applied",
        "reasoning written in plain language",
    ], 2.35, 4.25, 8.6, 1.35, fill=LIGHT_BLUE, accent=BLUE)


def add_semantic_slide(prs):
    slide = blank_slide(prs, "Semantic matching: beyond exact keywords", "AI matching")
    card(slide, "Why it exists", [
        "Bidder evidence may use different wording from tender clauses.",
        "Keyword-only matching misses equivalents like revenue vs turnover.",
        "Semantic score is shown instead of hidden.",
    ], 0.8, 1.25, 3.65, 2.65, fill=WHITE, accent=PURPLE)
    card(slide, "Implementation", [
        "SentenceTransformer('all-MiniLM-L6-v2') when installed",
        "normalize_embeddings=True",
        "Cosine similarity from embedding dot product",
        "Fallback uses token overlap + sequence similarity",
    ], 4.85, 1.25, 3.65, 2.65, fill=WHITE, accent=BLUE)
    card(slide, "Decision thresholds", [
        ">= 0.85 -> PASS",
        "0.72 to 0.85 -> NEEDS_REVIEW",
        "< 0.72 -> FAIL",
        "Missing evidence -> NEEDS_REVIEW",
    ], 8.9, 1.25, 3.0, 2.65, fill=WHITE, accent=TEAL)
    add_text(slide, "Prototype note", 0.95, 4.55, 2.1, 0.28, Pt(14), NAVY, bold=True)
    add_text(slide, "The local demo can run without downloading MiniLM. Installing requirements.txt enables the full sentence-transformers path.", 0.95, 4.95, 10.2, 0.45, Pt(13), BLACK)


def add_confidence_slide(prs):
    slide = blank_slide(prs, "Confidence logic: no silent rejection", "Decision policy")
    add_text(slide, "Semantic similarity thresholds", 0.8, 1.2, 4.2, 0.3, Pt(15), NAVY, bold=True)
    bar_x, bar_y, bar_w, bar_h = 0.85, 1.8, 10.8, 0.55
    segments = [
        (0.0, 0.72, RED, "FAIL"),
        (0.72, 0.85, AMBER, "NEEDS_REVIEW"),
        (0.85, 1.0, GREEN, "PASS"),
    ]
    for start, end, color, label in segments:
        x = bar_x + bar_w * start
        w = bar_w * (end - start)
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(bar_y), Inches(w), Inches(bar_h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        add_text(slide, label, x, bar_y + 0.17, w, 0.18, Pt(9), WHITE, bold=True, align=PP_ALIGN.CENTER)
    for threshold in [0.72, 0.85]:
        x = bar_x + bar_w * threshold
        connector(slide, x, bar_y - 0.15, x, bar_y + 0.8, NAVY)
        add_text(slide, str(threshold), x - 0.18, bar_y + 0.85, 0.45, 0.2, Pt(9), NAVY, bold=True)

    card(slide, "Rule path", [
        "If evidence is missing: NEEDS_REVIEW",
        "If evidence is confident and rule matches: PASS",
        "If evidence is confident and rule fails: FAIL",
        "If evidence confidence is below 0.85: NEEDS_REVIEW",
    ], 0.9, 3.55, 5.4, 1.85, fill=WHITE, accent=TEAL)
    card(slide, "Flagging rule", [
        "flagged=True for NEEDS_REVIEW",
        "flagged=True when evidence is missing",
        "flagged=True when confidence < 0.85",
        "UI surfaces these cases for officer validation",
    ], 6.85, 3.55, 5.1, 1.85, fill=WHITE, accent=AMBER)


def add_hitl_slide(prs):
    slide = blank_slide(prs, "Human-in-the-loop design", "Governance")
    card(slide, "System decides", [
        "What criteria exist",
        "What evidence was extracted",
        "Which rule was applied",
        "Confidence and recommendation",
    ], 0.75, 1.25, 3.3, 2.25, fill=WHITE, accent=BLUE)
    card(slide, "System refuses to decide alone", [
        "Missing documents",
        "Low-confidence OCR",
        "Ambiguous semantic score",
        "Unsafe parsing or unsupported rule",
    ], 4.45, 1.25, 3.3, 2.25, fill=WHITE, accent=AMBER)
    card(slide, "Officer reviews", [
        "Reason",
        "Source document",
        "Confidence score",
        "Extracted value and excerpt",
    ], 8.15, 1.25, 3.3, 2.25, fill=WHITE, accent=GREEN)
    add_text(slide, "This keeps AI in the assistant role: it speeds up extraction and consistency, while procurement authority remains with the human evaluator.", 1.05, 4.45, 10.9, 0.75, Pt(17), NAVY, bold=True)


def add_reporting_audit_slide(prs):
    slide = blank_slide(prs, "Stage 4: Reporting and audit logging", "Output")
    card(slide, "JSON report", [
        "Full EvaluationResult object",
        "All criteria and bidders",
        "Every CriterionEvaluation",
        "Machine-readable for integrations",
    ], 0.75, 1.25, 3.35, 2.3, fill=WHITE, accent=BLUE)
    card(slide, "HTML report", [
        "Evaluation matrix",
        "Per-bidder summary cards",
        "Decision colors",
        "Plain-English explanations",
    ], 4.55, 1.25, 3.35, 2.3, fill=WHITE, accent=TEAL)
    card(slide, "SQLite audit log", [
        "upload",
        "extraction",
        "evaluation",
        "timestamp + metadata_json",
    ], 8.35, 1.25, 3.35, 2.3, fill=WHITE, accent=PURPLE)
    card(slide, "Why auditors care", [
        "The system can show what document was processed, when it was processed, what criteria were extracted, and what report was generated.",
        "Every verdict carries enough context to defend or revise the decision.",
    ], 1.2, 4.35, 10.5, 1.2, fill=LIGHT_GREEN, accent=GREEN)


def add_sample_result_slide(prs, result):
    slide = blank_slide(prs, "Sample run: tender plus three bidders", "Demo output")
    if not result:
        add_text(slide, "Run scripts/smoke_test.py to generate sample output.", 0.9, 1.4, 8.0, 0.4, Pt(18), RED, bold=True)
        return
    summary = result.get("summary", {})
    metrics = [
        ("Criteria", len(result.get("criteria", [])), BLUE),
        ("Bidders", len(result.get("bidders", [])), TEAL),
        ("Evaluations", len(result.get("evaluations", [])), PURPLE),
        ("Reports", "JSON + HTML", GREEN),
    ]
    for idx, (label, value, color) in enumerate(metrics):
        card(slide, label, str(value), 0.75 + idx * 3.0, 1.15, 2.45, 0.9, fill=WHITE, accent=color, title_size=Pt(11), body_size=Pt(18))
    rows = [["Bidder", "Score", "Pass", "Fail", "Review", "Recommendation"]]
    for bidder, item in summary.items():
        rows.append([
            bidder,
            f"{item['score_percent']}%",
            str(item["pass"]),
            str(item["fail"]),
            str(item["needs_review"]),
            item["recommendation"],
        ])
    add_table(slide, rows, 0.75, 2.75, 11.9, 2.2)
    add_text(slide, "Interpretation", 0.8, 5.45, 1.8, 0.25, Pt(13), NAVY, bold=True)
    add_text(slide, "The demo intentionally includes a strong bidder, a bidder with hard failures, and a bidder with a review-worthy missing/uncertain case so judges can see all decision states.", 0.8, 5.82, 11.3, 0.42, Pt(12), BLACK)


def add_explainability_slide(prs, result):
    slide = blank_slide(prs, "What one verdict contains", "Explainability")
    sample = None
    if result:
        evaluations = result.get("evaluations", [])
        sample = next((item for item in evaluations if item.get("flagged")), None) or (evaluations[0] if evaluations else None)
    if sample:
        rows = [
            ["Field", "Value"],
            ["Bidder", sample.get("bidder_name", "")],
            ["Criterion", sample.get("criterion_name", "")],
            ["Decision", sample.get("decision", "")],
            ["Confidence", str(sample.get("confidence", ""))],
            ["Extracted value", str(sample.get("extracted_value", ""))],
            ["Source document", str(sample.get("source_document", ""))],
            ["Rule applied", sample.get("rule_applied", "")],
        ]
        add_table(slide, rows, 0.75, 1.15, 6.1, 4.75, font_size=Pt(9))
        card(slide, "Reasoning", sample.get("reasoning", "No sample reasoning available."), 7.25, 1.15, 5.0, 2.15, fill=WHITE, accent=AMBER)
    else:
        add_text(slide, "No report sample found.", 0.9, 1.3, 5, 0.4, Pt(16), RED, bold=True)
    card(slide, "Explainability contract", [
        "Criterion",
        "Extracted value",
        "Source document name",
        "Rule applied",
        "Final decision",
        "Reasoning and confidence",
    ], 7.25, 3.65, 5.0, 2.05, fill=LIGHT_BLUE, accent=BLUE)


def add_demo_run_slide(prs):
    slide = blank_slide(prs, "How to run the demo locally", "Operations")
    card(slide, "Terminal 1: backend", [
        'cd "C:\\Users\\Shivansh Karanwal\\Documents\\Codex\\2026-05-04\\hi-can-you-write-code-for"',
        ".\\.venv\\Scripts\\python.exe -m uvicorn backend.main:app --reload --port 8000",
        "Open http://127.0.0.1:8000/docs",
    ], 0.75, 1.25, 5.9, 2.05, fill=WHITE, accent=BLUE, body_size=Pt(9))
    card(slide, "Terminal 2: frontend", [
        'cd "C:\\Users\\Shivansh Karanwal\\Documents\\Codex\\2026-05-04\\hi-can-you-write-code-for"',
        ".\\.venv\\Scripts\\python.exe -m streamlit run frontend\\streamlit_app.py",
        "Open http://127.0.0.1:8501",
    ], 6.95, 1.25, 5.65, 2.05, fill=WHITE, accent=TEAL, body_size=Pt(9))
    card(slide, "Files to upload", [
        "data/sample/sample_tender_ai_services.pdf",
        "data/sample/sample_bidders.zip",
        "Click Upload Tender, Upload Bidders, Extract Criteria, Evaluate",
    ], 0.75, 4.05, 5.9, 1.55, fill=LIGHT_GREEN, accent=GREEN)
    card(slide, "Quick verification", [
        ".\\.venv\\Scripts\\python.exe scripts\\generate_sample_data.py",
        ".\\.venv\\Scripts\\python.exe scripts\\smoke_test.py",
    ], 6.95, 4.05, 5.65, 1.55, fill=LIGHT_BLUE, accent=BLUE)


def add_roadmap_slide(prs):
    slide = blank_slide(prs, "Roadmap from prototype to production", "Next steps")
    card(slide, "Short term", [
        "Officer approval/edit screen for extracted criteria",
        "Inline evidence viewer with highlighted excerpts",
        "Better validation for dates, certificates, and financial years",
    ], 0.75, 1.25, 3.55, 2.4, fill=WHITE, accent=BLUE)
    card(slide, "Medium term", [
        "Authenticated roles and officer IDs",
        "Hash-chained audit log and digital signatures",
        "PostgreSQL and object storage",
        "Background workers for large bidder batches",
    ], 4.85, 1.25, 3.55, 2.4, fill=WHITE, accent=TEAL)
    card(slide, "AI hardening", [
        "LLM provider integration with structured output validation",
        "Domain-specific embeddings and glossary",
        "Multilingual OCR and translation support",
        "Evaluation regression test set",
    ], 8.95, 1.25, 3.55, 2.4, fill=WHITE, accent=PURPLE)
    card(slide, "Production principle", "AI recommends, explains, and flags. Final procurement decision stays auditable and reviewable by authorized officers.", 1.25, 4.55, 10.8, 0.95, fill=LIGHT_AMBER, accent=AMBER)


def add_closing_slide(prs):
    slide = blank_slide(prs, "Closing argument", "Summary", dark=True)
    add_text(slide, "TenderMind AI converts a fragile manual workflow into a structured, explainable, auditable evaluation pipeline.", 0.85, 1.35, 11.2, 0.95, Pt(28), WHITE, bold=True)
    card(slide, "What judges should remember", [
        "End-to-end runnable prototype, not just a pitch.",
        "Clear service boundaries and typed schemas.",
        "PASS / FAIL / NEEDS_REVIEW for every criterion.",
        "Ambiguity is surfaced to humans, not hidden.",
        "Reports and audit logs make decisions defendable.",
    ], 0.95, 3.05, 5.6, 2.65, fill=RGBColor(238, 243, 250), accent=TEAL)
    card(slide, "Demo promise", "Upload tender, upload bidder ZIP, extract criteria, evaluate, inspect flagged cases, open the report.", 7.1, 3.05, 4.9, 1.55, fill=RGBColor(238, 243, 250), accent=GREEN)
    add_footer(slide, dark=True)


def add_table(slide, rows, x, y, w, h, font_size=Pt(9.5)):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    for r_idx, row in enumerate(rows):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(value)
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.fill.solid()
            if r_idx == 0:
                cell.fill.fore_color.rgb = NAVY
                color = WHITE
                bold = True
            else:
                decision_text = str(value)
                if "QUALIFIED" in decision_text and "NOT" not in decision_text:
                    cell.fill.fore_color.rgb = LIGHT_GREEN
                elif "NOT_QUALIFIED" in decision_text:
                    cell.fill.fore_color.rgb = LIGHT_RED
                elif "REVIEW" in decision_text:
                    cell.fill.fore_color.rgb = LIGHT_AMBER
                else:
                    cell.fill.fore_color.rgb = WHITE
                color = BLACK
                bold = False
            for p in cell.text_frame.paragraphs:
                p.font.size = font_size
                p.font.name = "Aptos"
                p.font.bold = bold
                p.font.color.rgb = color
    return table


if __name__ == "__main__":
    main()
