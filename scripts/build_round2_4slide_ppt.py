from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "TenderMind_AI_Round2_4_Slides.pptx"

NAVY = RGBColor(18, 32, 51)
TEAL = RGBColor(27, 132, 121)
BLUE = RGBColor(39, 96, 172)
GREEN = RGBColor(29, 116, 72)
RED = RGBColor(170, 52, 45)
AMBER = RGBColor(184, 123, 22)
PURPLE = RGBColor(102, 68, 167)
SLATE = RGBColor(74, 85, 104)
LIGHT = RGBColor(247, 249, 252)
WHITE = RGBColor(255, 255, 255)
BORDER = RGBColor(215, 222, 232)
LIGHT_BLUE = RGBColor(228, 238, 252)
LIGHT_GREEN = RGBColor(221, 247, 231)
LIGHT_RED = RGBColor(255, 225, 223)
LIGHT_AMBER = RGBColor(255, 242, 202)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_round2(prs)
    slide_dataset_making(prs)
    slide_sample_scenario(prs)
    slide_references(prs)

    prs.save(OUTPUT)
    print(f"Created {OUTPUT}")


def bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.18))
    band.fill.solid()
    band.fill.fore_color.rgb = TEAL
    band.line.fill.background()


def text(slide, value, x, y, w, h, size=14, color=NAVY, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = value
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    if align:
        p.alignment = align
    return box


def bullets(slide, items, x, y, w, h, size=12, color=SLATE):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"- {item}"
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(5)
    return box


def card(slide, title, body, x, y, w, h, accent=TEAL, fill=WHITE, title_size=13, body_size=11):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(1)

    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    text(slide, title, x + 0.18, y + 0.13, w - 0.32, 0.25, title_size, NAVY, True)
    if isinstance(body, list):
        bullets(slide, body, x + 0.2, y + 0.52, w - 0.35, h - 0.58, body_size, SLATE)
    else:
        text(slide, body, x + 0.18, y + 0.52, w - 0.35, h - 0.6, body_size, SLATE)
    return shape


def pill(slide, label, x, y, w, color, font_size=9):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.34))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    text(slide, label, x, y + 0.08, w, 0.16, font_size, WHITE, True, PP_ALIGN.CENTER)


def arrow(slide, x1, y1, x2, y2, color=SLATE):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(1.5)
    return line


def footer(slide):
    text(slide, "TenderMind AI | Round 2 prototype update", 0.55, 7.14, 5.2, 0.18, 8, SLATE)


def heading(slide, title, subtitle):
    text(slide, title, 0.55, 0.5, 11.6, 0.46, 25, NAVY, True)
    text(slide, subtitle, 0.58, 1.05, 11.4, 0.32, 13, SLATE)


def slide_round2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    heading(slide, "1. What is done in Round 2", "Updated from a basic prototype into a fuller sandbox-ready evaluation flow.")

    metrics = [
        ("12", "criteria extracted", BLUE),
        ("4", "bidder profiles", TEAL),
        ("48", "criterion-level verdicts", PURPLE),
        ("3", "decision states", GREEN),
    ]
    for idx, (num, label, color) in enumerate(metrics):
        card(slide, label, num, 0.7 + idx * 3.05, 1.68, 2.45, 0.88, color, WHITE, 11, 20)

    card(slide, "Round 2 implementation upgrades", [
        "Mandatory + optional criteria handling.",
        "Numeric, financial, date, boolean, text/semantic criteria.",
        "DOCX technical proposal parsing added.",
        "Image/photo OCR path included for scanned evidence.",
        "Optional failure no longer disqualifies overall bidder.",
    ], 0.75, 3.05, 3.75, 2.55, BLUE)

    card(slide, "AI evaluation logic", [
        "Tender -> Pydantic Criterion schema.",
        "Bidder docs -> normalized ExtractedField values.",
        "Rules evaluate hard criteria; semantic matcher checks technical proposal alignment.",
        "PASS / FAIL / NEEDS_REVIEW generated per criterion.",
    ], 4.85, 3.05, 3.75, 2.55, TEAL)

    card(slide, "Officer-focused output", [
        "Evaluation matrix across all bidders and criteria.",
        "Per-bidder recommendation: QUALIFIED / NOT_QUALIFIED / REVIEW_REQUIRED.",
        "HTML + JSON reports.",
        "SQLite audit logs for upload, extraction, evaluation.",
    ], 8.95, 3.05, 3.55, 2.55, GREEN)

    text(slide, "Core principle: the AI assists evaluation, but ambiguous or missing evidence is routed to human review instead of silent rejection.", 0.95, 6.18, 11.2, 0.35, 16, NAVY, True, PP_ALIGN.CENTER)
    footer(slide)


def slide_dataset_making(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    heading(slide, "2. Sample data set making", "Representative mock/redacted-style documents built to cover all important tender evaluation data types.")

    card(slide, "Tender PDF generated", [
        "AI-enabled Citizen Service Analytics Platform.",
        "Mandatory criteria: experience, turnover, projects, GST, ISO, dates, blacklisting, net worth, bid security, data residency, technical proposal.",
        "Optional criterion: local language support.",
    ], 0.75, 1.55, 3.7, 2.25, BLUE)

    card(slide, "Bidder ZIP generated", [
        "4 bidder folders inside sample_bidders.zip.",
        "Each bidder has mixed evidence files.",
        "Clearly eligible, clearly ineligible, optional-fail, and manual-review cases are included.",
    ], 4.82, 1.55, 3.7, 2.25, TEAL)

    card(slide, "Document formats covered", [
        "Typed PDF: company profile and certificates.",
        "TXT: compliance declaration.",
        "DOCX: technical proposal.",
        "PNG/JPG: scanned/photo evidence path.",
    ], 8.9, 1.55, 3.45, 2.25, PURPLE)

    rows = [
        ["Data Type", "Demo Criterion", "Evidence Source"],
        ["Number", "Experience >= 3 yrs, Projects >= 2", "company_profile.pdf"],
        ["Money", "Turnover >= INR 50 lakh", "company_profile.pdf"],
        ["Date", "ISO / Bid security valid until 2026-12-31", "PDF/TXT/JPG"],
        ["Boolean", "GST, ISO, blacklist, net worth, data residency", "TXT/PDF/image"],
        ["Text/Semantic", "Cloud-native AI platform alignment", "technical_proposal.docx"],
        ["Optional", "Local language support", "TXT/PDF declaration"],
    ]
    table(slide, rows, 0.8, 4.35, 11.8, 1.85, font_size=8.5)
    footer(slide)


def slide_sample_scenario(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    heading(slide, "3. Sample data scenario", "The demo intentionally creates clean pass, hard fail, optional fail, and manual-review behavior.")

    rows = [
        ["Bidder", "Scenario", "Key Evidence / Gap", "Outcome"],
        ["BharatTech Solutions", "Strong bidder", "All mandatory + optional evidence present", "QUALIFIED"],
        ["CivicAI Labs", "Optional fail only", "Mandatory pass; no local language support", "QUALIFIED"],
        ["NewWave Analytics", "Hard failures", "2 yrs exp, 45 lakh turnover, expired ISO, outside India", "NOT_QUALIFIED"],
        ["RuralLogic Innovations", "Review-heavy case", "Image-only/missing GST, ISO validity, bid security evidence", "NOT_QUALIFIED + REVIEW"],
    ]
    table(slide, rows, 0.75, 1.65, 11.85, 2.0, font_size=8.5)

    card(slide, "What judges can show live", [
        "Upload sample_tender_ai_services.pdf.",
        "Upload sample_bidders.zip.",
        "Click Extract Criteria -> 12 criteria appear.",
        "Click Evaluate -> 48 verdicts are produced.",
        "Open HTML report for matrix + explanations.",
    ], 0.85, 4.35, 3.55, 1.65, BLUE, WHITE, 12, 10)

    card(slide, "Expected decision behavior", [
        "Clearly eligible bidders become QUALIFIED.",
        "Mandatory failures become NOT_QUALIFIED.",
        "Optional failure affects score but not eligibility.",
        "Missing/low-confidence evidence becomes NEEDS_REVIEW.",
    ], 4.9, 4.35, 3.55, 1.65, TEAL, WHITE, 12, 10)

    card(slide, "Files used", [
        "data/sample/sample_tender_ai_services.pdf",
        "data/sample/sample_bidders.zip",
        "Generated by scripts/generate_sample_data.py",
    ], 8.95, 4.35, 3.35, 1.65, GREEN, WHITE, 12, 9.5)

    text(slide, "Demo result after update: 12 criteria | 4 bidders | 48 criterion-level decisions.", 1.0, 6.35, 11.2, 0.32, 16, NAVY, True, PP_ALIGN.CENTER)
    footer(slide)


def slide_references(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    heading(slide, "4. References: comparable systems, models, and technical stack", "Used as design references and technical grounding for the Round 2 prototype.")

    left = [
        ("GeM India", "https://gem.gov.in"),
        ("CPPP / eProcure India", "https://eprocure.gov.in/eprocure/app"),
        ("EU TED", "https://ted.europa.eu/en/about-ted"),
        ("SAM.gov opportunities", "https://sam.gov/content/opportunities"),
        ("World Bank STEP", "https://www.worldbank.org/en/projects-operations/products-and-services/brief/systematic-tracking-of-exchanges-in-procurement-step"),
        ("OCDS", "https://standard.open-contracting.org/latest/en/"),
    ]
    mid = [
        ("Google Document AI", "https://docs.cloud.google.com/document-ai/docs"),
        ("LayoutLM paper", "https://www.microsoft.com/en-us/research/publication/layoutlm-pre-training-of-text-and-layout-for-document-image-understanding/"),
        ("LayoutLMv3 paper", "https://arxiv.org/abs/2204.08387"),
        ("MiniLM model", "https://huggingface.co/sentence-transformers/all-MiniLM-L6-v2"),
        ("Tesseract OCR", "https://tesseract-ocr.github.io/"),
        ("OpenCV thresholding", "https://docs.opencv.org/4.x/d7/d4d/tutorial_py_thresholding.html"),
    ]
    right = [
        ("FastAPI", "https://fastapi.tiangolo.com/"),
        ("Streamlit", "https://docs.streamlit.io/"),
        ("Pydantic", "https://docs.pydantic.dev/"),
        ("pdfplumber", "https://github.com/jsvine/pdfplumber"),
        ("PyMuPDF", "https://pymupdf.readthedocs.io/"),
        ("SQLite", "https://www.sqlite.org/docs.html"),
    ]

    reference_block(slide, "Similar procurement / transparency systems", left, 0.65, 1.48, 4.0, BLUE)
    reference_block(slide, "Document AI / OCR / NLP models", mid, 4.82, 1.48, 4.0, PURPLE)
    reference_block(slide, "Prototype technical references", right, 8.98, 1.48, 3.7, TEAL)

    text(slide, "How these references map to TenderMind AI", 0.75, 6.05, 3.6, 0.22, 12, NAVY, True)
    text(
        slide,
        "Procurement portals inspired the upload/search/report workflow; OCDS inspired structured contracting data; LayoutLM/Document AI informed scanned-document direction; MiniLM/Tesseract/OpenCV support semantic and OCR paths; FastAPI/Streamlit/Pydantic/SQLite power the prototype.",
        0.75,
        6.38,
        11.7,
        0.45,
        10.5,
        SLATE,
    )
    footer(slide)


def reference_block(slide, title, refs, x, y, w, accent):
    card(slide, title, "", x, y, w, 4.35, accent, WHITE, 11, 9)
    y0 = y + 0.55
    for idx, (name, url) in enumerate(refs):
        text(slide, f"{name}", x + 0.18, y0 + idx * 0.58, w - 0.25, 0.16, 8.7, NAVY, True)
        text(slide, url, x + 0.18, y0 + 0.2 + idx * 0.58, w - 0.25, 0.17, 6.2, SLATE)


def table(slide, rows, x, y, w, h, font_size=8.5):
    shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = shape.table
    for r_idx, row in enumerate(rows):
        for c_idx, value in enumerate(row):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = str(value)
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.fill.solid()
            if r_idx == 0:
                cell.fill.fore_color.rgb = NAVY
                font_color = WHITE
                bold = True
            else:
                v = str(value)
                if "QUALIFIED" in v and "NOT" not in v:
                    cell.fill.fore_color.rgb = LIGHT_GREEN
                elif "NOT_QUALIFIED" in v:
                    cell.fill.fore_color.rgb = LIGHT_RED
                elif "REVIEW" in v:
                    cell.fill.fore_color.rgb = LIGHT_AMBER
                else:
                    cell.fill.fore_color.rgb = WHITE
                font_color = NAVY
                bold = False
            for p in cell.text_frame.paragraphs:
                p.font.name = "Aptos"
                p.font.size = Pt(font_size)
                p.font.bold = bold
                p.font.color.rgb = font_color
    return tbl


if __name__ == "__main__":
    main()

