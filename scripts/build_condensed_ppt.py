import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "TenderMind_AI_Condensed_3_Slides.pptx"

NAVY = RGBColor(18, 32, 51)
TEAL = RGBColor(27, 132, 121)
BLUE = RGBColor(39, 96, 172)
GREEN = RGBColor(29, 116, 72)
RED = RGBColor(170, 52, 45)
AMBER = RGBColor(184, 123, 22)
PURPLE = RGBColor(102, 68, 167)
SLATE = RGBColor(74, 85, 104)
LIGHT = RGBColor(247, 249, 252)
WHITE = RGBColor(255, 255, 255)
BORDER = RGBColor(215, 222, 232)


def main():
    result = latest_result()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_problem_solution(prs)
    slide_architecture(prs)
    slide_demo_result(prs, result)

    prs.save(OUTPUT)
    print(f"Created {OUTPUT}")


def latest_result():
    reports = sorted((ROOT / "data" / "reports").glob("eval_*.json"), key=lambda p: p.stat().st_mtime)
    if not reports:
        return None
    return json.loads(reports[-1].read_text(encoding="utf-8"))


def bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.18))
    band.fill.solid()
    band.fill.fore_color.rgb = TEAL
    band.line.fill.background()


def text(slide, value, x, y, w, h, size=14, color=NAVY, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = value
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    if align:
        p.alignment = align
    return box


def bullets(slide, items, x, y, w, h, size=13, color=SLATE):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"- {item}"
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return box


def card(slide, title, body, x, y, w, h, accent=TEAL, fill=WHITE, title_size=14, body_size=12):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(1)

    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    text(slide, title, x + 0.2, y + 0.14, w - 0.35, 0.26, title_size, NAVY, True)
    if isinstance(body, list):
        bullets(slide, body, x + 0.22, y + 0.55, w - 0.4, h - 0.65, body_size, SLATE)
    else:
        text(slide, body, x + 0.2, y + 0.58, w - 0.4, h - 0.7, body_size, SLATE)
    return shape


def pill(slide, label, x, y, w, color, font_size=10):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.38))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    text(slide, label, x, y + 0.09, w, 0.16, font_size, WHITE, True, PP_ALIGN.CENTER)
    return shape


def arrow(slide, x1, y1, x2, y2, color=SLATE):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(1.6)
    try:
        line.line.end_arrowhead = True
    except Exception:
        pass
    return line


def footer(slide):
    text(slide, "TenderMind AI | Explainable AI-based tender evaluation", 0.55, 7.14, 6.3, 0.18, 8, SLATE)


def slide_problem_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    text(slide, "TenderMind AI: AI-based Tender Evaluation System", 0.55, 0.55, 11.2, 0.45, 26, NAVY, True)
    text(slide, "A runnable prototype that converts unstructured tender and bidder documents into explainable, auditable eligibility decisions.", 0.58, 1.08, 11.4, 0.38, 14, SLATE)

    card(slide, "Problem", [
        "Tender criteria are scattered across long PDFs.",
        "Bidder documents arrive as typed PDFs, scans, images, and mixed formats.",
        "Manual checks are slow, inconsistent, and hard to defend later.",
    ], 0.7, 1.85, 3.55, 2.6, RED)

    card(slide, "Core promise", [
        "Extract eligibility criteria automatically.",
        "Parse bidder evidence with confidence scores.",
        "Return PASS / FAIL / NEEDS_REVIEW per criterion.",
        "Never silently reject ambiguous evidence.",
    ], 4.85, 1.85, 3.55, 2.6, TEAL)

    card(slide, "Why it matters", [
        "Officer gets a structured checklist.",
        "Every verdict has source, value, rule, confidence, and reasoning.",
        "Audit logs capture upload, extraction, and evaluation events.",
    ], 9.0, 1.85, 3.55, 2.6, BLUE)

    text(slide, "Outcome: faster evaluation, fewer missed criteria, clearer human review, and a report that can be audited.", 1.0, 5.35, 11.3, 0.5, 20, NAVY, True, PP_ALIGN.CENTER)
    footer(slide)


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    text(slide, "How the working prototype operates", 0.55, 0.55, 10.6, 0.45, 26, NAVY, True)
    text(slide, "FastAPI backend + Streamlit frontend + modular AI services + SQLite audit trail.", 0.58, 1.08, 11.0, 0.32, 14, SLATE)

    stages = [
        ("1", "Criteria Extraction", "Tender PDF -> Pydantic Criterion JSON\nLLM abstraction with mock fallback", BLUE),
        ("2", "Document Processing", "Bidder ZIP -> parsed PDFs/images\npdfplumber/PyMuPDF + OCR fallback", TEAL),
        ("3", "Matching Engine", "Rules + semantic similarity\nnumeric, boolean, date, confidence logic", PURPLE),
        ("4", "Reporting + Audit", "Matrix + explanations\nJSON/HTML report + SQLite logs", GREEN),
    ]
    x_positions = [0.65, 3.8, 6.95, 10.1]
    for idx, (num, title, body, color) in enumerate(stages):
        pill(slide, f"Stage {num}", x_positions[idx] + 0.15, 1.88, 1.0, color)
        card(slide, title, body, x_positions[idx], 2.35, 2.6, 2.05, color, WHITE, 13, 10.5)
        if idx < len(stages) - 1:
            arrow(slide, x_positions[idx] + 2.68, 3.35, x_positions[idx + 1] - 0.08, 3.35)

    card(slide, "Main API endpoints", [
        "POST /upload-tender",
        "POST /upload-bidders",
        "POST /extract-criteria",
        "POST /evaluate",
        "GET /report",
        "GET /audit-logs",
    ], 0.8, 5.0, 3.5, 1.45, BLUE, WHITE, 13, 9.5)

    card(slide, "Core code modules", [
        "criteria_extractor.py",
        "document_parser.py",
        "matcher.py",
        "rule_engine.py",
        "confidence_engine.py",
        "report_service.py",
    ], 4.95, 5.0, 3.5, 1.45, TEAL, WHITE, 13, 9.5)

    card(slide, "Decision policy", [
        ">= 0.85: PASS",
        "0.72-0.85: NEEDS_REVIEW",
        "< 0.72: FAIL",
        "Missing evidence: NEEDS_REVIEW",
    ], 9.1, 5.0, 3.4, 1.45, AMBER, WHITE, 13, 9.5)
    footer(slide)


def slide_demo_result(prs, result):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    text(slide, "Demo output: what judges will see", 0.55, 0.55, 10.6, 0.45, 26, NAVY, True)
    text(slide, "Upload sample tender + bidder ZIP, extract criteria, evaluate, then inspect matrix, flagged cases, and report.", 0.58, 1.08, 11.6, 0.32, 14, SLATE)

    if result:
        criteria = len(result.get("criteria", []))
        bidders = len(result.get("bidders", []))
        evaluations = len(result.get("evaluations", []))
        metrics = [("Criteria", criteria, BLUE), ("Bidders", bidders, TEAL), ("Verdicts", evaluations, PURPLE), ("Reports", "JSON + HTML", GREEN)]
    else:
        metrics = [("Criteria", "7", BLUE), ("Bidders", "3", TEAL), ("Verdicts", "21", PURPLE), ("Reports", "JSON + HTML", GREEN)]

    for idx, (label, value, color) in enumerate(metrics):
        card(slide, label, str(value), 0.7 + idx * 3.05, 1.75, 2.45, 0.95, color, WHITE, 12, 20)

    rows = [["Bidder", "Pass", "Fail", "Review", "Recommendation"]]
    summary = result.get("summary", {}) if result else {
        "BharatTech Solutions": {"pass": 7, "fail": 0, "needs_review": 0, "recommendation": "QUALIFIED"},
        "NewWave Analytics": {"pass": 5, "fail": 2, "needs_review": 0, "recommendation": "NOT_QUALIFIED"},
        "RuralLogic Innovations": {"pass": 5, "fail": 1, "needs_review": 1, "recommendation": "NOT_QUALIFIED"},
    }
    for bidder, item in summary.items():
        rows.append([bidder, item["pass"], item["fail"], item["needs_review"], item["recommendation"]])
    table(slide, rows, 0.8, 3.25, 7.4, 2.1)

    card(slide, "One verdict includes", [
        "Criterion name",
        "Extracted value",
        "Source document",
        "Rule applied",
        "Confidence score",
        "Plain-English reasoning",
    ], 8.75, 3.05, 3.55, 2.3, TEAL, WHITE, 13, 10)

    text(slide, "Demo files:", 0.9, 6.1, 1.3, 0.2, 12, NAVY, True)
    text(slide, "data/sample/sample_tender_ai_services.pdf    |    data/sample/sample_bidders.zip", 2.0, 6.1, 8.8, 0.2, 11, SLATE)
    text(slide, "Final takeaway: AI accelerates and structures evaluation, but ambiguous cases stay with the human evaluator.", 0.9, 6.55, 11.3, 0.32, 16, NAVY, True, PP_ALIGN.CENTER)
    footer(slide)


def table(slide, rows, x, y, w, h):
    shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = shape.table
    for r_idx, row in enumerate(rows):
        for c_idx, value in enumerate(row):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = str(value)
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.fill.solid()
            if r_idx == 0:
                cell.fill.fore_color.rgb = NAVY
                font_color = WHITE
                bold = True
            else:
                if str(value) == "QUALIFIED":
                    cell.fill.fore_color.rgb = RGBColor(221, 247, 231)
                elif "NOT_QUALIFIED" in str(value):
                    cell.fill.fore_color.rgb = RGBColor(255, 225, 223)
                elif "REVIEW" in str(value):
                    cell.fill.fore_color.rgb = RGBColor(255, 242, 202)
                else:
                    cell.fill.fore_color.rgb = WHITE
                font_color = NAVY
                bold = False
            for p in cell.text_frame.paragraphs:
                p.font.name = "Aptos"
                p.font.size = Pt(9)
                p.font.bold = bold
                p.font.color.rgb = font_color
    return tbl


if __name__ == "__main__":
    main()

